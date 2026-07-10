"""PowerPoint Presentation Parser using python-pptx."""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation

from src.observability.logging import get_logger

logger = get_logger(__name__)


class PptxParser:
    """Parses text placeholders and shapes inside PowerPoint presentation slides."""

    @staticmethod
    def parse(file_path: Path | str) -> str:
        """Extracts text content from .pptx file."""
        logger.info("Parsing PPTX file", path=str(file_path))
        try:
            prs = Presentation(file_path)
            extracted_text = []

            for i, slide in enumerate(prs.slides):
                slide_text = []
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    extracted_text.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_text))

            return "\n\n".join(extracted_text).strip()
        except Exception as e:
            logger.error("PPTX parsing failed", error=str(e), path=str(file_path))
            return ""
